"""
Serviço de exportação de aulas.

Formatos suportados:
  - PDF      : reportlab (layout com branding AprendAI)
  - Markdown : texto puro estruturado
  - PPTX     : python-pptx (um slide por seção)
  - CSV      : flashcards frente/verso (para importar no Anki, etc.)

Módulos opcionais (professor):
  - quiz         : perguntas + gabarito
  - flashcards   : cartões de memorização
  - methodology  : metodologias de ensino sugeridas pela IA
"""
import csv
import io
import logging
from dataclasses import dataclass
from typing import Literal

from reportlab.lib import colors
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import cm
from reportlab.platypus import (
    HRFlowable, PageBreak, Paragraph,
    SimpleDocTemplate, Spacer, Table, TableStyle,
)
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from app.models.schemas import LessonResponse, QuizResponse, FlashcardsResponse
from app.prompts.agents import get_methodology_system_prompt, get_methodology_user_prompt
from app.services.llm_client import call_model

logger = logging.getLogger(__name__)

# ─── Paleta AprendAI ──────────────────────────────────────────────────────────

PRIMARY     = colors.HexColor("#c8f060")   # verde limão
DARK_BG     = colors.HexColor("#0a0a08")
DARK_CARD   = colors.HexColor("#111110")
TEXT_LIGHT  = colors.HexColor("#e8e8e2")
TEXT_MUTED  = colors.HexColor("#6b6b60")

PRIMARY_RGB = RGBColor(0xC8, 0xF0, 0x60)
DARK_RGB    = RGBColor(0x0A, 0x0A, 0x08)
WHITE_RGB   = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_RGB    = RGBColor(0x6B, 0x6B, 0x60)


# ─── Dataclass de opções de exportação ────────────────────────────────────────

@dataclass
class ExportOptions:
    format: Literal["pdf", "markdown", "pptx", "csv"]
    include_quiz: bool = False
    include_flashcards: bool = False
    include_methodology: bool = False   # exclusivo professor
    # Dados pré-gerados (None = não incluir)
    quiz: QuizResponse | None = None
    flashcards: FlashcardsResponse | None = None
    methodology: dict | None = None


# ─── Agente de metodologias ───────────────────────────────────────────────────

async def generate_methodology(
    subject: str,
    lesson_title: str,
    level: str,
    tone: str,
) -> dict:
    logger.info(f"[ExportService] Gerando metodologias para: {lesson_title}")
    return await call_model(
        system_prompt=get_methodology_system_prompt(),
        user_prompt=get_methodology_user_prompt(subject, lesson_title, level, tone),
        max_tokens=2000,
        temperature=0.5,
    )


# ─── PDF ──────────────────────────────────────────────────────────────────────

def export_pdf(lesson: LessonResponse, opts: ExportOptions) -> bytes:
    buffer = io.BytesIO()
    doc = SimpleDocTemplate(
        buffer, pagesize=A4,
        leftMargin=2*cm, rightMargin=2*cm,
        topMargin=2*cm, bottomMargin=2*cm,
    )

    styles = getSampleStyleSheet()

    # Estilos customizados
    title_style = ParagraphStyle(
        "AprendTitle",
        parent=styles["Title"],
        fontSize=22, leading=28,
        textColor=DARK_BG, spaceAfter=6,
    )
    h2_style = ParagraphStyle(
        "AprendH2",
        parent=styles["Heading2"],
        fontSize=14, leading=18,
        textColor=DARK_BG, spaceBefore=16, spaceAfter=6,
    )
    body_style = ParagraphStyle(
        "AprendBody",
        parent=styles["Normal"],
        fontSize=11, leading=16,
        textColor=colors.HexColor("#222220"), spaceAfter=8,
    )
    label_style = ParagraphStyle(
        "AprendLabel",
        parent=styles["Normal"],
        fontSize=8, leading=10,
        textColor=TEXT_MUTED,
        fontName="Helvetica",
        spaceAfter=4,
    )
    italic_style = ParagraphStyle(
        "AprendItalic",
        parent=styles["Normal"],
        fontSize=11, leading=16,
        textColor=colors.HexColor("#333330"),
        fontName="Helvetica-Oblique",
        leftIndent=16, spaceAfter=8,
    )

    story = []

    # ── Cabeçalho ──────────────────────────────────────────────────────────────
    story.append(Paragraph("APRENDAI", label_style))
    story.append(Paragraph(lesson.title, title_style))
    story.append(Paragraph(
        f"Aula {lesson.lesson_number} &nbsp;·&nbsp; {lesson.estimated_reading_minutes} min de leitura",
        label_style,
    ))
    story.append(HRFlowable(width="100%", thickness=2, color=PRIMARY, spaceAfter=16))

    # ── Seções ─────────────────────────────────────────────────────────────────
    for section in lesson.sections:
        story.append(Paragraph(section.heading, h2_style))
        story.append(Paragraph(section.body.replace("\n", "<br/>"), body_style))

    # ── Conceitos-chave ────────────────────────────────────────────────────────
    if lesson.key_concepts:
        story.append(Spacer(1, 12))
        story.append(HRFlowable(width="100%", thickness=1, color=colors.HexColor("#dddddd")))
        story.append(Paragraph("CONCEITOS-CHAVE", label_style))
        for concept in lesson.key_concepts:
            story.append(Paragraph(f"• {concept}", body_style))

    # ── Reflexão ───────────────────────────────────────────────────────────────
    if lesson.reflection_question:
        story.append(Spacer(1, 8))
        story.append(Paragraph("PERGUNTA DE REFLEXÃO", label_style))
        story.append(Paragraph(lesson.reflection_question, italic_style))

    # ── Quiz ───────────────────────────────────────────────────────────────────
    if opts.include_quiz and opts.quiz:
        story.append(PageBreak())
        story.append(Paragraph("QUIZ DE AVALIAÇÃO", label_style))
        story.append(HRFlowable(width="100%", thickness=2, color=PRIMARY, spaceAfter=12))

        for q in opts.quiz.questions:
            story.append(Paragraph(f"{q.number}. {q.question}", h2_style))
            for opt in q.options:
                marker = "✓" if opt.letter == q.correct_letter else "   "
                color  = colors.HexColor("#1a6a00") if opt.letter == q.correct_letter else colors.HexColor("#444442")
                story.append(Paragraph(
                    f"{marker} {opt.letter}) {opt.text}",
                    ParagraphStyle("opt", parent=body_style, textColor=color),
                ))
            story.append(Paragraph(f"<i>Explicação: {q.explanation}</i>", italic_style))
            story.append(Spacer(1, 8))

    # ── Flashcards ─────────────────────────────────────────────────────────────
    if opts.include_flashcards and opts.flashcards:
        story.append(PageBreak())
        story.append(Paragraph("FLASHCARDS", label_style))
        story.append(HRFlowable(width="100%", thickness=2, color=PRIMARY, spaceAfter=12))

        data = [["#", "Frente", "Verso"]] + [
            [str(i+1), card.front, card.back]
            for i, card in enumerate(opts.flashcards.cards)
        ]
        table = Table(data, colWidths=[1*cm, 8*cm, 8*cm])
        table.setStyle(TableStyle([
            ("BACKGROUND",  (0, 0), (-1, 0), PRIMARY),
            ("TEXTCOLOR",   (0, 0), (-1, 0), DARK_BG),
            ("FONTNAME",    (0, 0), (-1, 0), "Helvetica-Bold"),
            ("FONTSIZE",    (0, 0), (-1, 0), 9),
            ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#f5f5f3")]),
            ("FONTSIZE",    (0, 1), (-1, -1), 9),
            ("VALIGN",      (0, 0), (-1, -1), "TOP"),
            ("GRID",        (0, 0), (-1, -1), 0.5, colors.HexColor("#dddddd")),
            ("PADDING",     (0, 0), (-1, -1), 6),
        ]))
        story.append(table)

    # ── Metodologias (professor) ───────────────────────────────────────────────
    if opts.include_methodology and opts.methodology:
        story.append(PageBreak())
        story.append(Paragraph("METODOLOGIAS DE ENSINO", label_style))
        story.append(HRFlowable(width="100%", thickness=2, color=PRIMARY, spaceAfter=12))

        for m in opts.methodology.get("methodologies", []):
            story.append(Paragraph(m["name"], h2_style))
            story.append(Paragraph(
                f"<b>Categoria:</b> {m.get('category','')} &nbsp;·&nbsp; "
                f"<b>Tempo estimado:</b> {m.get('estimated_minutes','')} min",
                label_style,
            ))
            story.append(Paragraph(m.get("description", ""), body_style))
            if m.get("steps"):
                story.append(Paragraph("<b>Passo a passo:</b>", body_style))
                for step in m["steps"]:
                    story.append(Paragraph(f"• {step}", body_style))
            if m.get("materials"):
                mats = ", ".join(m["materials"])
                story.append(Paragraph(f"<b>Materiais:</b> {mats}", body_style))
            story.append(Spacer(1, 12))

        # Objetivos de aprendizagem
        if opts.methodology.get("learning_objectives"):
            story.append(Paragraph("OBJETIVOS DE APRENDIZAGEM", label_style))
            for obj in opts.methodology["learning_objectives"]:
                story.append(Paragraph(f"• {obj}", body_style))

        # Sugestões de avaliação
        if opts.methodology.get("assessment_suggestions"):
            story.append(Spacer(1, 8))
            story.append(Paragraph("SUGESTÕES DE AVALIAÇÃO", label_style))
            for sug in opts.methodology["assessment_suggestions"]:
                story.append(Paragraph(f"• {sug}", body_style))

    doc.build(story)
    return buffer.getvalue()


# ─── Markdown ─────────────────────────────────────────────────────────────────

def export_markdown(lesson: LessonResponse, opts: ExportOptions) -> bytes:
    lines = []

    lines.append(f"# {lesson.title}")
    lines.append(f"\n> Aula {lesson.lesson_number} · {lesson.estimated_reading_minutes} min de leitura\n")

    for section in lesson.sections:
        lines.append(f"\n## {section.heading}\n")
        lines.append(section.body)

    if lesson.key_concepts:
        lines.append("\n## Conceitos-chave\n")
        for concept in lesson.key_concepts:
            lines.append(f"- {concept}")

    if lesson.reflection_question:
        lines.append(f"\n## Reflexão\n\n> {lesson.reflection_question}")

    if opts.include_quiz and opts.quiz:
        lines.append("\n---\n\n## Quiz\n")
        for q in opts.quiz.questions:
            lines.append(f"\n**{q.number}. {q.question}**\n")
            for opt in q.options:
                marker = "✓" if opt.letter == q.correct_letter else " "
                lines.append(f"- [{marker}] {opt.letter}) {opt.text}")
            lines.append(f"\n*Explicação: {q.explanation}*")

    if opts.include_flashcards and opts.flashcards:
        lines.append("\n---\n\n## Flashcards\n")
        lines.append("| # | Frente | Verso |")
        lines.append("|---|--------|-------|")
        for i, card in enumerate(opts.flashcards.cards):
            lines.append(f"| {i+1} | {card.front} | {card.back} |")

    if opts.include_methodology and opts.methodology:
        lines.append("\n---\n\n## Metodologias de Ensino\n")
        for m in opts.methodology.get("methodologies", []):
            lines.append(f"\n### {m['name']}")
            lines.append(f"**Categoria:** {m.get('category','')} · **Tempo:** {m.get('estimated_minutes','')} min\n")
            lines.append(m.get("description", ""))
            if m.get("steps"):
                lines.append("\n**Passo a passo:**")
                for step in m["steps"]:
                    lines.append(f"1. {step}")

    lines.append(f"\n\n---\n*Gerado por AprendAI*")
    return "\n".join(lines).encode("utf-8")


# ─── CSV (flashcards) ─────────────────────────────────────────────────────────

def export_csv(lesson: LessonResponse, opts: ExportOptions) -> bytes:
    buffer = io.StringIO()
    writer = csv.writer(buffer)
    writer.writerow(["frente", "verso", "aula"])

    # Conceitos-chave como flashcards básicos
    for concept in lesson.key_concepts:
        writer.writerow([f"O que é {concept}?", concept, lesson.title])

    # Flashcards gerados pela IA
    if opts.include_flashcards and opts.flashcards:
        for card in opts.flashcards.cards:
            writer.writerow([card.front, card.back, lesson.title])

    # Quiz como flashcards pergunta/resposta
    if opts.include_quiz and opts.quiz:
        for q in opts.quiz.questions:
            correct = next((o.text for o in q.options if o.letter == q.correct_letter), "")
            writer.writerow([q.question, correct, lesson.title])

    return buffer.getvalue().encode("utf-8")


# ─── PPTX ─────────────────────────────────────────────────────────────────────

def export_pptx(lesson: LessonResponse, opts: ExportOptions) -> bytes:
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # Totalmente em branco

    def _bg(slide):
        """Aplica fundo escuro em todos os slides."""
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = DARK_RGB

    def _add_label(slide, text, top):
        txb = slide.shapes.add_textbox(Inches(0.6), top, Inches(6), Inches(0.3))
        tf  = txb.text_frame
        tf.text = text.upper()
        p = tf.paragraphs[0]
        p.runs[0].font.size = Pt(8)
        p.runs[0].font.color.rgb = GRAY_RGB
        p.runs[0].font.bold = True

    def _add_title(slide, text, top=Inches(0.5)):
        txb = slide.shapes.add_textbox(Inches(0.6), top, Inches(12), Inches(1.2))
        tf  = txb.text_frame
        tf.word_wrap = True
        tf.text = text
        p = tf.paragraphs[0]
        p.runs[0].font.size = Pt(32)
        p.runs[0].font.bold = True
        p.runs[0].font.color.rgb = WHITE_RGB

    def _add_body(slide, text, top, height=Inches(4.5)):
        txb = slide.shapes.add_textbox(Inches(0.6), top, Inches(12), height)
        tf  = txb.text_frame
        tf.word_wrap = True
        tf.text = text
        p = tf.paragraphs[0]
        p.runs[0].font.size = Pt(14)
        p.runs[0].font.color.rgb = RGBColor(0xC8, 0xC8, 0xC0)

    def _accent_bar(slide):
        """Linha verde limão horizontal no topo."""
        bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(0.06))
        bar.fill.solid()
        bar.fill.fore_color.rgb = PRIMARY_RGB
        bar.line.fill.background()

    # ── Slide 1: Capa ──────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    _bg(slide)

    # Faixa primária lateral
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.3), Inches(7.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY_RGB
    bar.line.fill.background()

    _add_label(slide, f"AprendAI · Aula {lesson.lesson_number}", Inches(1.2))
    _add_title(slide, lesson.title, Inches(2.5))
    _add_label(slide, f"{lesson.estimated_reading_minutes} min de leitura", Inches(4.2))

    # ── Slides de seções ───────────────────────────────────────────────────────
    for section in lesson.sections:
        slide = prs.slides.add_slide(blank_layout)
        _bg(slide)
        _accent_bar(slide)
        _add_label(slide, section.heading, Inches(0.5))
        _add_body(slide, section.body, Inches(0.9))

    # ── Slide: Conceitos-chave ─────────────────────────────────────────────────
    if lesson.key_concepts:
        slide = prs.slides.add_slide(blank_layout)
        _bg(slide)
        _accent_bar(slide)
        _add_label(slide, "Conceitos-chave", Inches(0.5))

        txb = slide.shapes.add_textbox(Inches(0.6), Inches(0.9), Inches(12), Inches(5.5))
        tf  = txb.text_frame
        tf.word_wrap = True
        for i, concept in enumerate(lesson.key_concepts):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            p.text = f"• {concept}"
            p.runs[0].font.size = Pt(16)
            p.runs[0].font.color.rgb = PRIMARY_RGB
            p.space_after = Pt(8)

    # ── Slide: Quiz ────────────────────────────────────────────────────────────
    if opts.include_quiz and opts.quiz:
        for q in opts.quiz.questions:
            slide = prs.slides.add_slide(blank_layout)
            _bg(slide)
            _accent_bar(slide)
            _add_label(slide, f"Quiz · Pergunta {q.number}", Inches(0.5))

            txb = slide.shapes.add_textbox(Inches(0.6), Inches(0.9), Inches(12), Inches(1.5))
            tf  = txb.text_frame
            tf.word_wrap = True
            tf.text = q.question
            tf.paragraphs[0].runs[0].font.size = Pt(18)
            tf.paragraphs[0].runs[0].font.color.rgb = WHITE_RGB

            opt_txb = slide.shapes.add_textbox(Inches(0.6), Inches(2.6), Inches(12), Inches(3.5))
            opt_tf  = opt_txb.text_frame
            opt_tf.word_wrap = True
            for i, opt in enumerate(q.options):
                is_correct = opt.letter == q.correct_letter
                p = opt_tf.add_paragraph() if i > 0 else opt_tf.paragraphs[0]
                p.text = f"{'✓' if is_correct else '  '} {opt.letter}) {opt.text}"
                p.runs[0].font.size = Pt(13)
                p.runs[0].font.color.rgb = PRIMARY_RGB if is_correct else RGBColor(0xC8, 0xC8, 0xC0)
                p.space_after = Pt(6)

    # ── Slide: Metodologias ────────────────────────────────────────────────────
    if opts.include_methodology and opts.methodology:
        for m in opts.methodology.get("methodologies", []):
            slide = prs.slides.add_slide(blank_layout)
            _bg(slide)
            _accent_bar(slide)
            _add_label(slide, f"Metodologia · {m.get('category', '')}", Inches(0.5))
            _add_label(slide, m["name"], Inches(0.85))

            body  = m.get("description", "")
            steps = m.get("steps", [])
            if steps:
                body += "\n\n" + "\n".join(f"• {s}" for s in steps[:4])
            _add_body(slide, body, Inches(1.2), Inches(5))

    # ── Slide final ────────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    _bg(slide)
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY_RGB
    bar.line.fill.background()
    _add_title(slide, "AprendAI", Inches(3.0))
    _add_label(slide, "Motor de Aprendizado com Inteligência Artificial", Inches(4.4))

    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()